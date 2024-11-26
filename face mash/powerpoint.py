from pptx import Presentation
from pptx.util import Pt
from arabic_reshaper import reshape
from bidi.algorithm import get_display

def add_rtl_text(frame, text, font_name="B Nazanin", font_size=24):
    """افزودن متن راست‌چین و فارسی به فریم متن در اسلاید"""
    reshaped_text = reshape(text)
    bidi_text = get_display(reshaped_text)
    frame.text = bidi_text
    for paragraph in frame.text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)

# ایجاد ارائه پاورپوینت
ppt = Presentation()

# اسلاید 1: عنوان
slide = ppt.slides.add_slide(ppt.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
add_rtl_text(title, "سیستم شناسایی و تحلیل چهره", font_size=32)
add_rtl_text(subtitle, "با استفاده از پایتون، OpenCV، MediaPipe و face_recognition", font_size=20)

# اسلاید 2: معرفی سیستم
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
add_rtl_text(title, "معرفی سیستم")
add_rtl_text(
    content,
    "- شناسایی و تحلیل چهره در زمان واقعی\n"
    "- استفاده از یادگیری ماشین و بینایی کامپیوتر\n"
    "- اجزای کلیدی: ثبت چهره، شناسایی، و تشخیص بر اساس عکس"
)

# اسلاید 3: کتابخانه‌های مورد استفاده
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
add_rtl_text(title, "کتابخانه‌های مورد استفاده")
add_rtl_text(
    content,
    "- OpenCV: برای دریافت ویدئو و ترسیم المان‌های بصری\n"
    "- MediaPipe: برای شناسایی نقاط صورت\n"
    "- face_recognition: برای کدگذاری و تطبیق چهره‌ها\n"
    "- NumPy: برای عملیات عددی\n"
    "- pickle: برای ذخیره و بارگذاری کدگذاری چهره‌ها\n"
    "- time: برای عملیات مبتنی بر زمان"
)

# اسلاید 4: فرآیند ثبت چهره
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
add_rtl_text(title, "فرآیند ثبت چهره")
add_rtl_text(
    content,
    "- دریافت ویدئو از وب‌کم\n"
    "- استفاده از MediaPipe برای شناسایی نقاط چهره\n"
    "- تولید کدگذاری چهره با face_recognition\n"
    "- ذخیره کدگذاری میانگین برای شناسایی در آینده\n"
    "- راهنمایی کاربر برای چرخاندن سر جهت پوشش بهتر"
)

# اسلاید 5: شناسایی با ویدئو
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
add_rtl_text(title, "شناسایی با ویدئو")
add_rtl_text(
    content,
    "- دریافت ویدئو به صورت لحظه‌ای\n"
    "- بارگذاری کدگذاری چهره ذخیره‌شده\n"
    "- مقایسه کدگذاری چهره شناسایی‌شده با کدگذاری ذخیره‌شده\n"
    "- نمایش وضعیت تطبیق با بازخورد بصری\n"
    "- مکانیزم محدودیت زمانی برای افزایش کارایی"
)

# اسلاید 6: شناسایی با عکس
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
add_rtl_text(title, "شناسایی با عکس")
add_rtl_text(
    content,
    "- امکان گرفتن عکس توسط کاربر برای شناسایی\n"
    "- استفاده از face_recognition برای کدگذاری چهره\n"
    "- مقایسه کدگذاری با داده‌های ذخیره‌شده\n"
    "- ارائه وضعیت تطبیق با بازخورد بصری"
)

# اسلاید 7: ویژگی‌های کلیدی
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
add_rtl_text(title, "ویژگی‌های کلیدی")
add_rtl_text(
    content,
    "- تحلیل چهره در زمان واقعی\n"
    "- پشتیبانی از شناسایی با ویدئو و عکس\n"
    "- ارائه راهنمای بصری و بازخورد پویا\n"
    "- مدیریت سناریوهای مختلف شناسایی\n"
    "- ارائه تعامل کاربرپسند"
)

# اسلاید 8: مراحل الگوریتم
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
add_rtl_text(title, "مراحل الگوریتم")
add_rtl_text(
    content,
    "1. دریافت ورودی ویدئو یا عکس\n"
    "2. شناسایی چهره با MediaPipe\n"
    "3. کدگذاری چهره با face_recognition\n"
    "4. مقایسه کدگذاری با داده‌های ذخیره‌شده\n"
    "5. نمایش نتایج با فاصله و وضعیت تطبیق"
)

# اسلاید 9: ساختار کد
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
add_rtl_text(title, "ساختار کد")
add_rtl_text(
    content,
    "- capture_and_store_encodings(): ثبت چهره\n"
    "- recognize_face_with_video(): شناسایی چهره با ویدئو\n"
    "- recognize_face_with_photo(): شناسایی چهره با عکس\n"
    "- توابع ماژولار و قابل استفاده مجدد"
)

# اسلاید 10: نتیجه‌گیری و کاربردها
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
add_rtl_text(title, "نتیجه‌گیری و کاربردها")
add_rtl_text(
    content,
    "- ترکیب کتابخانه‌های پیشرفته برای شناسایی کارآمد چهره\n"
    "- کاربردهای بالقوه در امنیت، حضور و غیاب، و شخصی‌سازی\n"
    "- نمایش قدرت پایتون در حل مسائل دنیای واقعی"
)

# ذخیره فایل پاورپوینت
file_path = "/mnt/data/Face_Recognition_System_Presentation_FA.pptx"
ppt.save(file_path)

file_path
